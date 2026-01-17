import os, sqlite3
from pathlib import Path
from pypdf import PdfReader
from pptx import Presentation
import docx2txt
from ebooklib import epub
from .util import chunk_text
from .emb import embed_texts
from .store import _connect

def _detect_mime(path:str)->str:
    ext = Path(path).suffix.lower()
    return {
      ".pdf":"application/pdf",
      ".pptx":"application/vnd.openxmlformats-officedocument.presentationml.presentation",
      ".docx":"application/vnd.openxmlformats-officedocument.wordprocessingml.document",
      ".epub":"application/epub+zip"
    }.get(ext,"application/octet-stream")

def _read_pdf(p):
    r = PdfReader(p)
    for i, page in enumerate(r.pages, start=1):
        yield i, (page.extract_text() or "")

def _read_pptx(p):
    pr = Presentation(p)
    for i, s in enumerate(pr.slides, start=1):
        txt = "\n".join([sh.text for sh in s.shapes if hasattr(sh,"text")])
        yield i, txt

def _read_docx(p):
    yield 1, (docx2txt.process(p) or "")

def _read_epub(p):
    book = epub.read_epub(p)
    items = [it for it in book.get_items() if it.get_type()==9]
    html = " ".join([it.get_body_content().decode("utf-8",errors="ignore") for it in items])
    import re; yield 1, re.sub("<[^>]+>", " ", html)

def ingest_files(paths:list[str], db_path:str, index_dir:str, emb_model:str):
    os.makedirs(index_dir, exist_ok=True)
    con = _connect(db_path)
    try:
        cur = con.cursor()
        stats = []
        for p in paths:
            mime = _detect_mime(p)
            reader = { "application/pdf":_read_pdf, 
                       "application/vnd.openxmlformats-officedocument.presentationml.presentation":_read_pptx,
                       "application/vnd.openxmlformats-officedocument.wordprocessingml.document":_read_docx,
                       "application/epub+zip":_read_epub }.get(mime)
            if not reader:
                continue
            pages = list(reader(p))
            cur.execute(
                "INSERT INTO sources(filename,mime,pages,imported_at) VALUES(?,?,?,datetime('now'))",
                (os.path.basename(p), mime, len(pages)),
            )
            sid = cur.lastrowid
            chunks, payload = [], []
            for page, full in pages:
                for ch in chunk_text(full, max_chars=1100, overlap=200):
                    quote = (ch[:180] + "…") if len(ch) > 180 else ch
                    chunks.append((sid, page, ch, quote))
                    payload.append(ch)
            if payload:
                embs = embed_texts(payload, model_name=emb_model)
                for (sid, page, ch, quote), emb in zip(chunks, embs):
                    cur.execute(
                        "INSERT INTO chunks(source_id,page,text,quote,embedding) VALUES(?,?,?,?,?)",
                        (sid, page, ch, quote, emb),
                    )
            stats.append({"file": os.path.basename(p), "chunks": len(payload)})
        con.commit()
        return stats
    finally:
        con.close()

